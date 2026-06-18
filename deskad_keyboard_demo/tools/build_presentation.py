#!/usr/bin/env python
"""DeskAd 발표 산출물 빌더 — 프로젝트 보고서(PDF) + 슬라이드(PPTX).

자산:
  docs/presentation/assets/img/*.png   예시 이미지(UI·3D·depth·광고)
  docs/presentation/assets/fonts/NanumGothic-*.ttf  한글 폰트(없으면 download_font로 받음)
출력:
  docs/presentation/DeskAd_프로젝트_보고서.pdf
  docs/presentation/DeskAd_발표.pptx

사용: conda run -n sprint_high python tools/build_presentation.py
모든 docs/* 는 기본 .gitignore로 로컬 전용(공개 정책). 인프라 값은 본문에서 마스킹.
"""
from __future__ import annotations

import urllib.request
from pathlib import Path

from PIL import Image as PILImage

ROOT = Path(__file__).resolve().parent.parent
PRES = ROOT / "docs" / "presentation"
IMG = PRES / "assets" / "img"
FONTS = PRES / "assets" / "fonts"
PDF_OUT = PRES / "DeskAd_프로젝트_보고서.pdf"
PPTX_OUT = PRES / "DeskAd_발표.pptx"

FONT_URLS = {
    "NanumGothic-Regular.ttf": "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumgothic/NanumGothic-Regular.ttf",
    "NanumGothic-Bold.ttf": "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumgothic/NanumGothic-Bold.ttf",
}


def ensure_fonts() -> tuple[Path, Path]:
    FONTS.mkdir(parents=True, exist_ok=True)
    for name, url in FONT_URLS.items():
        p = FONTS / name
        if not p.exists():
            print("downloading font", name)
            urllib.request.urlretrieve(url, p)
    return FONTS / "NanumGothic-Regular.ttf", FONTS / "NanumGothic-Bold.ttf"


def img_size(name: str) -> tuple[int, int] | None:
    p = IMG / name
    if not p.exists():
        return None
    with PILImage.open(p) as im:
        return im.size


# ── 색상 팔레트 ───────────────────────────────────────────────────────────────
INK = "#1f2430"
ACCENT = "#c8552e"
MUTED = "#5a616b"
PANEL = "#f4f1eb"
BLUE = "#3766c4"


# ============================================================================
# 1) PDF 보고서 (reportlab Platypus)
# ============================================================================
def build_pdf() -> None:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    reg, bold = ensure_fonts()
    pdfmetrics.registerFont(TTFont("Nanum", str(reg)))
    pdfmetrics.registerFont(TTFont("Nanum-Bold", str(bold)))
    pdfmetrics.registerFontFamily("Nanum", normal="Nanum", bold="Nanum-Bold")

    from reportlab.platypus import (
        Image,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    styles = getSampleStyleSheet()

    def st(name, **kw):
        base = dict(fontName="Nanum", textColor=colors.HexColor(INK), leading=16)
        base.update(kw)
        return ParagraphStyle(name, parent=styles["Normal"], **base)

    body = st("body", fontSize=10.2, leading=16, spaceAfter=6)
    h1 = st("h1", fontName="Nanum-Bold", fontSize=18, textColor=colors.HexColor(ACCENT),
            spaceBefore=14, spaceAfter=8, leading=22)
    h2 = st("h2", fontName="Nanum-Bold", fontSize=13, textColor=colors.HexColor(INK),
            spaceBefore=10, spaceAfter=5, leading=17)
    cap = st("cap", fontSize=8.4, textColor=colors.HexColor(MUTED), alignment=TA_CENTER,
             spaceBefore=3, spaceAfter=10, leading=11)
    bullet = st("bullet", fontSize=10.2, leading=15, leftIndent=12, spaceAfter=3, bulletIndent=2)
    title = st("title", fontName="Nanum-Bold", fontSize=30, alignment=TA_LEFT,
               textColor=colors.HexColor(INK), leading=36)
    subtitle = st("subtitle", fontSize=13, textColor=colors.HexColor(MUTED), leading=20)

    story: list = []

    def para(text, style=body):
        story.append(Paragraph(text, style))

    def bullets(items):
        for it in items:
            story.append(Paragraph(f"•&nbsp;&nbsp;{it}", bullet))

    def figure(name, caption, max_w=15.5):
        sz = img_size(name)
        if not sz:
            return
        w, h = sz
        draw_w = min(max_w * cm, 16 * cm)
        draw_h = draw_w * h / w
        max_h = 11 * cm
        if draw_h > max_h:
            draw_h = max_h
            draw_w = draw_h * w / h
        story.append(Image(str(IMG / name), width=draw_w, height=draw_h))
        story.append(Paragraph(caption, cap))

    def figure_row(pairs, max_w=15.8):
        # pairs: [(name, caption), ...] 가로 배치
        cells, capcells = [], []
        cw = (max_w * cm) / len(pairs) - 0.25 * cm
        for name, caption in pairs:
            sz = img_size(name)
            if not sz:
                cells.append("")
                capcells.append("")
                continue
            w, h = sz
            dw = cw
            dh = dw * h / w
            cells.append(Image(str(IMG / name), width=dw, height=dh))
            capcells.append(Paragraph(caption, cap))
        t = Table([cells, capcells], colWidths=[cw + 0.25 * cm] * len(pairs))
        t.setStyle(TableStyle([("ALIGN", (0, 0), (-1, -1), "CENTER"),
                               ("VALIGN", (0, 0), (-1, -1), "TOP")]))
        story.append(t)
        story.append(Spacer(1, 0.3 * cm))

    def table(rows, header=True, col_widths=None):
        data = [[Paragraph(str(c), st("cell", fontSize=9, leading=12,
                                      fontName="Nanum-Bold" if (header and ri == 0) else "Nanum",
                                      textColor=colors.white if (header and ri == 0) else colors.HexColor(INK)))
                 for c in row] for ri, row in enumerate(rows)]
        t = Table(data, colWidths=col_widths, hAlign="LEFT")
        ts = [("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(INK)),
              ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#d8d3c8")),
              ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f7f5f0")]),
              ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
              ("TOPPADDING", (0, 0), (-1, -1), 4), ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
              ("LEFTPADDING", (0, 0), (-1, -1), 6)]
        t.setStyle(TableStyle(ts))
        story.append(t)
        story.append(Spacer(1, 0.35 * cm))

    # ── 표지 ──
    story.append(Spacer(1, 2.4 * cm))
    para("DeskAd AI Studio", title)
    story.append(Spacer(1, 0.3 * cm))
    para("3D 데스크 셋업 미리보기 + 배열 충실도 기반 광고 콘텐츠 자동 생성", subtitle)
    story.append(Spacer(1, 0.55 * cm))
    para("상세 프로젝트 보고서", st("ptag", fontName="Nanum-Bold", fontSize=15,
                          textColor=colors.HexColor(ACCENT)))
    story.append(Spacer(1, 0.2 * cm))
    para("작성일 2026-06-18 · 최종 업데이트 반영본 · 소상공인 커스텀 키보드/데스크테리어 판매자용", subtitle)
    story.append(Spacer(1, 0.7 * cm))
    figure("setup_3d_render.png", "절차적 3D 셋업 렌더 — 실제 상품 스펙과 데스크 구성을 광고 생성의 기준 데이터로 사용", max_w=13.5)
    para("본 보고서는 README, docs/architecture.md, docs/schema.md, docs/security.md, docs/deploy.md, 2026-06-10~15 작업 로그, "
         "그리고 backend/ui 핵심 구현을 바탕으로 작성했다. 외부 접속 주소·토큰·비밀번호 같은 민감 인프라 값은 포함하지 않는다.",
         st("note", fontSize=9, textColor=colors.HexColor(MUTED), leading=13))
    story.append(PageBreak())

    # ── 1. 요약 ──
    para("1. Executive Summary", h1)
    para("DeskAd AI Studio는 커스텀 키보드와 데스크테리어 판매자가 <b>상품 스펙 입력 → 정확한 3D 셋업 → 광고 문구 후보 → 실사 이미지 → "
         "SVG/PPTX 포스터</b>까지 한 화면 흐름으로 완성할 수 있게 만든 프로토타입이다. 프로젝트의 핵심 가치는 단순 이미지 생성이 아니라, "
         "판매자가 입력한 배열·색상·소품 배치를 광고 이미지가 얼마나 충실히 따르는지를 통제하는 데 있다.")
    para("최종 업데이트는 세 영역으로 정리된다. 첫째, Step 4 광고 콘텐츠 흐름을 문구·이미지·포스터 순서로 명확히 분리해 사용자가 실패 상태와 "
         "다음 행동을 이해하도록 했다. 둘째, local+ComfyUI 경로에서 depth-ControlNet과 grid_three 순차 생성으로 배열 충실도와 GPU 안정성을 높였다. "
         "셋째, 인증·보안·배포·캐시·품질 게이트를 보강해 데모 환경에서 반복 실행 가능한 운영 형태로 정리했다.")
    table([["구분", "최종 상태"],
           ["서비스 범위", "상품 정보, 도면/업로드, 3D 셋업, 광고 문구 후보, 이미지 job, 포스터 SVG/PPTX"],
           ["사용자 경험", "Streamlit 4단계 위저드, 로그인/회원가입 게이트, 작업 상태 카드, 실시간 진행 게이지"],
           ["핵심 기술", "절차적 GLB 렌더링, 셋업 구도 맵, GLB 기반 depth-ControlNet, best-of-N 액센트 색 선별"],
           ["AI 트랙", "OpenAI API 트랙과 Local+ComfyUI 트랙. 엔진 명시 시 다른 엔진으로 조용히 우회하지 않음"],
           ["운영/보안", "nginx Basic Auth, 세션 인증, 시크릿 마스킹, 0600 런타임 파일, GPU 워커 수명주기, Docker 앱 티어"],
           ["검증", "회귀 테스트와 라이브 ComfyUI 생성 검증을 병행. 시각 품질은 눈 검증으로 보완"]],
          col_widths=[3.4 * cm, 12.2 * cm])

    para("핵심 성과", h2)
    bullets([
        "<b>쓸 수 있는 광고 이미지</b>: 범용 text2image의 대표 실패인 풀사이즈 오생성, 키캡 녹음, 행 뒤틀림을 depth 구조 고정으로 줄였다.",
        "<b>상인 관점의 완성 흐름</b>: 포스터 생성까지 연결해 단일 이미지를 넘어서 실제 판매 채널에 올릴 산출물을 만든다.",
        "<b>운영 가능한 프로토타입</b>: API 키가 없어도 폴백으로 동작하고, GPU 워커·캐시·job store가 장시간 작업을 다룬다.",
        "<b>최종 업데이트 반영</b>: 로그인 UX, 작업 상태 안내, grid_three 시점 분리, 순차 큐잉, OpenAI/local 엔진 선택, Docker 앱 티어 전환 문서를 반영했다.",
    ])

    # ── 2. 배경 ──
    story.append(PageBreak())
    para("2. 배경과 문제 정의", h1)
    para("소상공인 판매자는 촬영 장비, 제품 스타일링, 광고 카피 작성, 상세페이지 편집에 투자할 여력이 제한적이다. 생성형 이미지 모델은 제작 속도를 "
         "크게 줄일 수 있지만, 커스텀 키보드처럼 배열·키 수·색 조합이 판매 상품 자체인 영역에서는 ‘그럴듯하지만 틀린 이미지’가 곧 상품 오표현이 된다.")
    table([["문제", "구체 현상", "사업 영향"],
           ["촬영/디자인 비용", "제품별 사진 촬영, 배경 연출, 상세페이지 작업이 반복됨", "신제품 등록 주기가 길어지고 테스트 판매가 어려움"],
           ["생성 이미지 불일치", "65% 배열이 풀배열처럼 보이거나, 키 열이 휘고 키캡이 녹음", "실제 상품과 다른 광고가 되어 신뢰를 해침"],
           ["도구 분산", "CAD/렌더러, LLM, 이미지 생성, 포스터 편집이 각각 다른 도구에 있음", "비전문 판매자가 끝까지 완성하기 어려움"],
           ["GPU/키 의존성", "모델 워커, API 키, 네트워크 상태에 따라 결과가 중단됨", "데모나 운영 중 실패 복구 경험이 나빠짐"]],
          col_widths=[3.0 * cm, 6.6 * cm, 6.0 * cm])
    para("DeskAd는 이 문제를 ‘광고 자동 생성’이 아니라 <b>상품 스펙을 기준 데이터로 삼는 생성 파이프라인</b>으로 접근했다. 입력한 배열·소품·치수로 "
         "먼저 3D 셋업을 만들고, 이 셋업을 이미지 생성 단계의 구조 레퍼런스로 사용한다. 최종 포스터 텍스트는 이미지에 직접 굽지 않고 SVG 레이어에서 "
         "합성해 한글 깨짐과 마케팅 문구 왜곡도 줄인다.")

    para("대상 사용자와 사용 시나리오", h2)
    bullets([
        "커스텀 키보드 공방 또는 키캡/데스크매트 판매자가 신제품별 광고 소재를 빠르게 제작한다.",
        "촬영 전 여러 데스크 배치와 색상 조합을 비교해 상세페이지 대표 컷 방향을 정한다.",
        "인스타그램, 스마트스토어, 상세페이지, 배너 광고 등 채널별 비율과 톤을 바꿔 소재를 반복 생성한다.",
        "디자이너가 없는 팀에서도 문구 후보를 비교하고 선택 문구를 직접 수정한 뒤 포스터로 내보낸다.",
    ])

    # ── 3. 최종 업데이트 ──
    story.append(PageBreak())
    para("3. 최종 업데이트 반영 사항", h1)
    para("최종 업데이트는 사용자 흐름 정리, 인증/세션 안정화, 생성 파이프라인 안정화, 운영 문서화로 나눌 수 있다. 2026-06-10~15 작업 로그와 "
         "현재 코드 상태를 기준으로 정리하면 다음과 같다.")
    table([["날짜/영역", "주요 변경", "효과"],
           ["2026-06-10 · Step 4 UX", "문구 생성 → 이미지 작업 → 포스터 생성 순서 카드, 버튼 활성 조건, 실패 후 재시도 상태 추가", "사용자가 다음 작업을 명확히 알 수 있고 실패 상태에서 멈추지 않음"],
           ["2026-06-10 · 상태 유지", "상품 입력 위젯 key와 저장용 session_state key 분리, 판매가 콤마 정규화", "단계 이동 후 상품 브리프가 초기화되는 문제 완화"],
           ["2026-06-11 · 오류 메시지", "API 연결 실패, timeout, HTTP 오류를 한국어 안내로 변환", "백엔드 미기동/설정 누락 상황을 사용자에게 설명 가능"],
           ["2026-06-11 · 로딩 최적화", "Step 2~4의 무거운 helper를 lazy import로 전환", "첫 화면 import 부담 감소, 로그인/초기 화면 진입 체감 개선"],
           ["2026-06-12 · 인증", "SignupResponse 분리, 로그인/회원가입 테스트 보강", "백엔드 인증 응답 모델과 테스트 범위 정리"],
           ["2026-06-15 · 로그인 UX", "브랜드 패널 + 로그인/회원가입 카드, URL 토큰 노출 방식 제거, session_state 기반 전환", "인증 화면 완성도와 보안성 개선"],
           ["2026-06-16~17 · 이미지 품질", "GLB depth-ControlNet, 컷별 depth 카메라, 유니크 depth 파일명, grid_three 순차 큐잉", "배열 충실도와 단일 L4 안정성 개선"],
           ["2026-06-15~배포", "Docker 앱 티어 문서와 컷오버 절차, host GPU 워커 연결 방식 정리", "FastAPI/Streamlit 앱 계층을 CPU-only 컨테이너로 운영 가능"]],
          col_widths=[3.3 * cm, 7.3 * cm, 5.0 * cm])
    para("특히 grid_three는 단일 프롬프트 batch로 처리하면 세 컷이 같은 시점의 노이즈 변형이 되기 쉽고, 한꺼번에 큐잉하면 FLUX+ControlNet이 단일 L4 VRAM을 "
         "압박한다. 현재 구현은 첫 컷만 큐에 올린 뒤 폴링이 완료를 확인할 때 다음 컷을 제출한다. 이로써 사용자에게는 3컷 광고가 보이고, 워커에는 항상 "
         "한 컷씩만 실행되는 구조가 된다.")

    # ── 4. 아키텍처 ──
    story.append(PageBreak())
    para("4. 시스템 아키텍처", h1)
    para("DeskAd는 외부에 노출되는 Streamlit UI와 내부 FastAPI 백엔드, 그리고 선택적으로 연결되는 텍스트/이미지 GPU 워커로 구성된다. 백엔드는 "
         "입력 검증, 렌더링, 프롬프트 구성, job store, cache, 워커 수명주기를 담당하는 단일 오케스트레이터다.")
    table([["계층", "주요 파일", "책임"],
           ["프론트엔드", "streamlit_app.py, ui/steps.py, ui/ad_content.py, ui/result_panel.py", "로그인 게이트, 4단계 위저드, 생성 버튼/상태 카드, 포스터 미리보기와 다운로드"],
           ["API 백엔드", "backend/main.py, backend/schemas.py, backend/app_factory.py", "REST 엔드포인트, Pydantic 입력 검증, 정적 파일 마운트, 헬스/보안 설정"],
           ["렌더링", "backend/renderer.py, backend/cad.py, backend/drawing_converter.py", "절차적 GLB 생성, STEP/GLB 업로드 처리, 플레이트 도면 변환, depth 렌더"],
           ["AI 생성", "backend/ai.py, backend/llm_adapters.py, backend/copy_policy.py", "문구 생성, 이미지 프롬프트, OpenAI/Local/ComfyUI 경로, 포스터 SVG"],
           ["운영 상태", "backend/job_store.py, backend/result_cache.py, backend/runtime_workers.py", "비동기 이미지 job, SHA256 캐시, GPU 워커 예열/해제/idle reap"],
           ["보안/인증", "backend/auth.py, backend/user_store.py, backend/security.py", "세션, 회원가입 코드, PBKDF2 사용자 저장, 시크릿 마스킹, 로그 필터"],
           ["배포", "Dockerfile, docker-compose.yml, start.sh, docs/deploy.md", "CPU-only 앱 티어, host GPU 워커 연결, nginx 경유 외부 접근"]],
          col_widths=[2.6 * cm, 5.7 * cm, 7.3 * cm])
    para("엔진 선택은 사용자-facing으로 <b>OpenAI API</b>와 <b>Local + ComfyUI</b> 두 트랙이다. OpenAI 트랙은 텍스트와 이미지를 OpenAI 호환 API로 "
         "처리한다. Local 트랙은 로컬 텍스트 후보(HyperCLOVA, Kanana, Mi:dm 등 설정된 provider)와 ComfyUI/FLUX 이미지 생성을 묶는다. "
         "엔진을 명시하면 평가 트랙의 무결성을 위해 다른 엔진으로 조용히 우회하지 않는다.")
    figure("ui_3d_preview.png", "Streamlit 3D 미리보기 화면 — 사용자가 생성한 GLB를 model-viewer로 확인", max_w=14)

    para("API 엔드포인트 요약", h2)
    table([["분류", "엔드포인트", "역할"],
           ["상태/보안", "GET /health, GET /security/config", "헬스체크와 마스킹된 설정 진단"],
           ["인증", "POST /auth/login, /auth/signup, /auth/logout, /auth/session", "세션 토큰 발급·검증·종료, 가입 코드 기반 회원가입"],
           ["3D 렌더", "POST /render/keyboard-preview, /render/desk-setup, /render/uploaded-model", "키보드 단품/데스크 셋업/업로드 모델 GLB 생성"],
           ["AI 문구", "POST /ai/copy, /ai/copy/variants, /ai/copy/experiment", "단일 문구, 후보 문구, provider 비교 생성"],
           ["AI 이미지", "POST /ai/image, POST /ai/image/jobs, GET /ai/image/jobs/{id}", "동기 이미지와 비동기 이미지 job 큐/폴링"],
           ["품질/포스터", "POST/GET /ai/image/jobs/{id}/quality, POST /ai/poster", "이미지 품질 평가와 SVG 포스터 합성"]],
          col_widths=[2.6 * cm, 5.7 * cm, 7.3 * cm])

    # ── 5. 4단계 UX ──
    story.append(PageBreak())
    para("5. 사용자 흐름과 UI 설계", h1)
    para("프론트엔드는 상품 등록자가 자연스럽게 순서를 따라가도록 4단계 위저드로 구성했다. 각 단계는 다음 단계의 입력 품질을 높이는 역할을 갖는다. "
         "예를 들어 Step 3에서 생성된 model_url과 composition_b64는 Step 4 이미지 프롬프트와 ControlNet 입력의 근거가 된다.")
    table([["단계", "사용자 입력", "구현 처리", "다음 단계에 넘기는 값"],
           ["1. 상품 정보", "상품 유형, 상품명, 판매가, 채널, 타깃, 핵심 특징, 상세 설명", "가격 숫자 정규화, 필수 입력 확인, widget/store state 분리", "카피 컨텍스트, 채널별 구도 기본값, 포스터 표시 텍스트"],
           ["2. 도면/제품 데이터", "키보드 모델/배열, STEP/STP/GLB, 공용 도면, 상세 커스텀", "레이아웃 JSON 조회, 업로드 모델 검증, 스위치/보강판/키캡 프로파일 저장", "배열/재질/색상/업로드 모델 메타"],
           ["3. 가상 셋업", "책상 크기, 모니터, 소품, 색상, 테마", "데스크 충돌 회피 배치, GLB 생성, model-viewer 표시, 구도 맵 생성", "model_url, composition_b64, composition_topdown_b64"],
           ["4. 광고 콘텐츠", "엔진, 톤, 비율, shot_type, 포스터 템플릿, 추가 요청", "문구 후보, 이미지 job, 자동 폴링, 포스터 생성/다운로드", "SVG/PPTX 포스터와 이미지 품질 리포트"]],
          col_widths=[2.7 * cm, 4.0 * cm, 5.0 * cm, 3.9 * cm])
    para("Step 4는 최종 업데이트에서 가장 크게 정리된 영역이다. 광고 문구가 생성되기 전에는 이미지 작업 버튼이 비활성화되고, 이미지 job이 완료되기 전에는 "
         "이미지 포함 포스터 생성이 비활성화된다. 이미지 job이 실패하거나 설정 누락으로 끝나면 ‘다시 실행 가능’ 상태가 표시되며, 사용자는 이미지 없이 "
         "SVG 포스터를 먼저 만들 수 있다.")
    para("로그인 화면도 중앙 카드형에서 브랜드 패널 + 인증 카드 구조로 바뀌었다. URL query parameter에 인증 토큰을 싣는 방식은 Streamlit 로그인 UI에서 "
         "사용하지 않고, 로그인 성공 후 session_state에 토큰을 저장해 즉시 메인 화면으로 전환한다. 기존 HttpOnly cookie 복원 경로는 남겨 두되 현재 UI는 "
         "자동 top-level 이동에 의존하지 않는다.")

    para("상태 관리와 UX 방어", h2)
    bullets([
        "Step 1 입력값은 위젯 key와 저장 key를 분리해, 다른 단계에서 위젯이 사라져도 상품 브리프가 초기화되지 않는다.",
        "광고 액션 notice는 rerun 뒤에도 한 번 표시되도록 session_state에 저장한다.",
        "이미지 job 폴링은 provider별 예상 시간을 사용하되 완료를 단정하지 않고 97%에서 멈춘다.",
        "grid_three는 실제 컷 완료 수를 기준으로 진행률을 표시해 시간 추정보다 신뢰할 수 있는 상태를 제공한다.",
        "운영자 진단 정보는 DESKAD_OPERATOR_MODE가 켜진 경우에만 표시해 일반 사용자 화면에서 모델/키 상태 노출을 줄인다.",
    ])

    # ── 6. 렌더링 ──
    story.append(PageBreak())
    para("6. 3D 렌더링과 셋업 데이터", h1)
    para("렌더링은 외부 CAD 렌더러에 전적으로 의존하지 않고, backend/renderer.py의 GlbBuilder가 박스·실린더·구체·토러스 primitives를 조립해 GLB를 "
         "생성한다. 이 접근은 프로토타입에서 모델 파일 준비 비용을 낮추고, 사용자가 바꾸는 배열·색상·소품·책상 치수를 즉시 반영하는 데 유리하다.")
    table([["구현 요소", "내용"],
           ["단위 규약", "1 GLB unit = 1 cm. MX 1u spacing = 1.905 cm. 모니터암은 VESA MIS-D 100 x 100 mm 기준"],
           ["키보드", "레이아웃 JSON의 key 좌표를 cm로 변환하고 case, plate, PCB, switch housing, keycap skirt/top을 계층적으로 생성"],
           ["커스텀 옵션", "case_finish, plate_material, pcb_color, switch_stem, switch_family, keycap_profile, mount_type, show_internals"],
           ["데스크", "책상 폭/깊이 clamp, 데스크매트 크기 자동 산정, 목재 결, 케이블 grommet, 접촉 그림자"],
           ["소품 배치", "DeskPlacer가 2D bounding box를 관리해 소품 겹침을 피하고, 후보 slot에서 가능한 위치를 선택"],
           ["출력 메타", "key_count, board_width/depth, case_outer_width/depth, monitor_panel_cm, placed_items, enabled_assets"],
           ["구도 맵", "실제 배치 좌표를 PIL로 원근/탑다운 PNG로 그려 img2img 레퍼런스로 전달"]],
          col_widths=[3.3 * cm, 12.3 * cm])
    figure_row([("composition_perspective.png", "원근 구도 맵 — hero/eye-level/wide 계열"),
                ("composition_topdown.png", "탑다운 구도 맵 — flat-lay 계열")])
    para("구도 맵은 최종 사진이 아니라 구조 신호다. DeskPlacer가 이미 알고 있는 책상 cm 좌표와 키보드 key footprint를 사용해 순수 PIL로 빠르게 그린다. "
         "top_down shot에서는 세워진 모니터와 충돌하지 않도록 탑다운 맵을 사용하고, hero/eye_level/wide_scene에서는 원근 맵 또는 depth 입력을 사용한다.")

    para("업로드 모델 처리", h2)
    bullets([
        "업로드 확장자는 STEP, STP, GLB만 허용한다.",
        "MAX_UPLOAD_MB 기본 60MB로 대용량 업로드에 의한 메모리/디스크 고갈을 제한한다.",
        "GLB는 glTF magic header를 검사하고, 바운딩 박스 기준으로 m/mm 단위 오해 가능성을 advisory로 알려 준다.",
        "STEP converter 명령이 없으면 proxy GLB를 만들어 model-viewer 흐름이 끊기지 않게 한다.",
        "파일명은 SHA256 일부와 timestamp 기반으로 재명명해 클라이언트 제공 경로를 그대로 신뢰하지 않는다.",
    ])

    # ── 7. AI 문구 ──
    story.append(PageBreak())
    para("7. 광고 문구 생성", h1)
    para("문구 생성은 backend/ai.py가 담당한다. 프론트는 Step 4에서 선택한 엔진과 사용자 입력을 build_ad_payload()로 묶어 넘기고, 백엔드는 "
         "selected_copy가 있으면 LLM을 우회해 사용자가 고른 문구를 그대로 포스터와 이미지 작업에 반영한다.")
    table([["구성", "구현"],
           ["Provider", "OpenAI, HyperCLOVA, Kanana, Mi:dm, local_llm, fallback. Local 트랙은 설정된 후보 provider를 비교한다."],
           ["프롬프트", "system prompt + tone별 few-shot + 상품/타깃/재질/배열 context. vision provider는 제품 레퍼런스 이미지를 첨부 가능"],
           ["인젝션 대응", "product_name, selling_point, product_detail, extra_request 등 자유텍스트에서 한/영 prompt injection 패턴을 flag-only로 기록"],
           ["후처리", "copy_policy가 과장/보장/순위/의학적 표현을 완화하고, 채널별 해시태그 수를 제한"],
           ["캐시", "provider, model, 상품 정보, 배열, 색상, 재질, 타깃 등을 포함한 SHA256 text cache key를 사용"],
           ["UX", "문구 후보는 provider/variant별 카드로 표시하고, 사용자가 선택/편집한 문구를 이후 이미지와 포스터 단계에서 사용"]],
          col_widths=[3.0 * cm, 12.6 * cm])
    para("fallback 문구는 단순 예외 처리가 아니라 데모 지속성을 위한 하한선이다. API 키나 로컬 LLM이 없어도 사용자는 포스터 생성 흐름까지 확인할 수 있다. "
         "다만 실제 품질 비교를 위해 엔진을 명시적으로 선택한 경우에는 다른 이미지 backend로 조용히 우회하지 않는 정책을 둔다.")

    para("카피 정책", h2)
    bullets([
        "‘국내 1위’, ‘최저가’, ‘100% 보장’, ‘완벽’, ‘치료’처럼 광고 리스크가 큰 표현을 완화 표현으로 치환한다.",
        "공백 회피 표현, 예를 들어 ‘국 내 1 위’처럼 음절 사이 공백을 넣은 경우도 정규식으로 탐지한다.",
        "정책 단계에서는 사용자가 편집한 긴 문구를 임의로 자르지 않고, 포스터 SVG 레이어에서 폭 기준 줄바꿈과 글자 크기 조정을 수행한다.",
    ])

    # ── 8. 이미지 생성 ──
    story.append(PageBreak())
    para("8. 이미지 생성 파이프라인", h1)
    para("이미지 프롬프트는 상품명과 카피를 단순히 붙이는 수준이 아니라, 구도·렌즈·조명·색온도·배열 제약·소품 인벤토리·네거티브 프롬프트를 함께 구성한다. "
         "이미지에는 마케팅 문구를 직접 렌더하지 말라고 명시하고, 텍스트는 포스터 SVG 단계에서 올린다.")
    table([["프롬프트 절", "목적"],
           ["subject", "layout, desk size, monitor size, material, selected assets를 피사체 정보로 명시"],
           ["keyboard fidelity", "정확한 키 수, 직선 행, crisp legends, accurate proportions를 양성 신호로 반복"],
           ["layout fidelity", "104는 numpad 포함, compact 배열은 numpad 금지처럼 배열별 제약을 분리"],
           ["exact colours", "case/keycap/accent/PCB 색을 피사체 직후에 배치해 색 드리프트를 줄임"],
           ["scene inventory", "선택한 소품을 exactly one of each로 명시해 마우스/소품 복제를 억제"],
           ["composition", "hero, top_down, detail_macro, eye_level, wide_scene별 angle/lens/framing 분기"],
           ["negative", "melted keycaps, warped rows, duplicate keyboard, two mice, random gadgets, gibberish text 금지"]],
          col_widths=[3.4 * cm, 12.2 * cm])
    para("이미지 backend는 OpenAI, HyperCLOVA image, local image endpoint, ComfyUI를 지원한다. 장시간 걸리는 경로는 /ai/image/jobs로 큐에 넣고 "
         "프론트가 3초마다 폴링한다. job에는 width/height, aspect_ratio, backend_config, prompt_preview, accent_keycap_color가 저장되어 품질 평가와 "
         "best-of-N 선별에 사용된다.")

    para("depth-ControlNet: 구조와 외관의 분리", h2)
    para("순수 text2image 또는 평면 구도 맵 img2img는 사진 품질과 정확한 키 배열을 동시에 얻기 어렵다. DeskAd는 Step 3에서 만든 GLB를 OSMesa "
         "소프트웨어 렌더로 depth PNG로 변환하고, ComfyUI FLUX depth-ControlNet에 넣어 배열과 소품 구조를 고정한다. 이 렌더는 GPU를 사용하지 않으므로 "
         "exclusive GPU 워커와 충돌하지 않는다.")
    figure_row([("composition_perspective.png", "① 셋업 구도"),
                ("depth_hero.png", "② GLB depth 입력"),
                ("ad_hero.png", "③ 사진 광고")])
    table([["환경 변수", "역할"],
           ["COMFYUI_CONTROLNET_MODEL", "ControlNet 모델 파일명. 비어 있으면 depth 경로 비활성"],
           ["COMFYUI_CONTROLNET_STRENGTH", "구조 강제 강도. 0보다 크고 모델이 설정되어야 ControlNet 후보 활성"],
           ["COMFYUI_CONTROLNET_END_PERCENT", "초기 스텝만 구조를 강제해 후반 사진감을 살리는 직교 노브"],
           ["COMFYUI_COMPOSITION_DENOISE", "구도 맵 img2img 전용 denoise. 셋업 색블록은 높은 denoise가 필요"],
           ["COMFYUI_BEST_OF_N", "N장 batch 중 액센트 색 충실도가 높은 후보를 quality_gate로 선택"]],
          col_widths=[7.0 * cm, 8.6 * cm])
    para("depth는 grayscale라 색을 직접 잠글 수 없다. 따라서 색상은 프롬프트의 exact colours 절과 best-of-N의 액센트 색 평가가 담당한다. 결과적으로 "
         "<b>배열은 depth, 주색은 프롬프트, 액센트는 best-of-N</b>으로 역할을 분리했다.")

    # ── 9. grid_three ──
    story.append(PageBreak())
    para("9. Grid Three와 비동기 Job 안정화", h1)
    para("grid_three 포스터는 제품 메인 컷, 키캡/스위치 디테일 컷, 데스크 무드 컷을 각각 다른 shot_type으로 생성한다. 기존 단일 batch 방식은 "
         "세 장이 같은 카메라의 변형으로 나오기 쉬웠고, depth-ControlNet을 동시에 큐잉하면 VRAM 피크가 겹칠 위험이 있었다.")
    figure_row([("depth_hero.png", "hero depth — 높은 3/4"),
                ("depth_eye_level.png", "eye_level depth — 낮은 수평"),
                ("depth_legacy.png", "legacy depth — 고정 각도")])
    para("현재 구현은 _grid_three_shot_plan이 hero/detail_macro/eye_level 세 컷을 정의하고, _submit_comfyui_grid_job이 첫 컷만 ComfyUI에 제출한다. "
         "프론트 폴링이 /history에서 완료를 확인하면 _poll_comfyui_grid_job이 다음 pending 컷을 제출한다. detail_macro는 카메라가 매크로 장면과 맞지 않으므로 "
         "ControlNet depth 경로를 쓰지 않고 img2img 또는 기본 워크플로로 폴백한다.")
    figure_row([("ad_hero.png", "최종 hero 컷"),
                ("ad_eye_level.png", "최종 eye_level 컷"),
                ("ad_detail_macro.png", "최종 detail_macro 컷")])
    table([["위험", "대응"],
           ["동시 큐잉으로 VRAM 피크 겹침", "ComfyUI 큐에 우리 컷이 항상 1개만 있도록 순차 제출"],
           ["hero/eye_level이 같은 depth 파일을 읽음", "shot_type과 depth content hash를 포함한 유니크 파일명 업로드"],
           ["best-of-N batch OOM", "ComfyUI history의 workflow에서 EmptyLatentImage batch_size를 반감해 재제출"],
           ["백엔드 재시작 후 running job 고착", "provider별 stale budget과 heartbeat로 failed 종결"],
           ["대용량 image_b64 응답", "목록 API와 public job에서 이미지 바이트를 마스킹하고 단건/포스터 경로에서만 사용"]],
          col_widths=[5.0 * cm, 10.6 * cm])

    # ── 10. 포스터 ──
    story.append(PageBreak())
    para("10. 포스터 생성과 산출물", h1)
    para("포스터 생성은 backend/ai.py의 create_svg_poster와 템플릿 빌더가 담당한다. 텍스트는 이미지에 직접 렌더하지 않고 SVG 텍스트 레이어로 올린다. "
         "이 방식은 한글이 이미지 생성 모델 안에서 깨지는 문제를 피하고, 사용자가 선택/수정한 문구를 포스터에 정확히 반영할 수 있게 한다.")
    table([["템플릿", "용도와 특징"],
           ["minimal_card", "이미지와 텍스트를 균형 있게 배치하는 기본 카드. CTA 대비를 검사해 버튼 색을 보정"],
           ["grid_three", "hero/detail/lifestyle 3컷을 메인·디테일·무드 패널에 배치. 1장만 있을 때도 다른 crop으로 fallback"],
           ["feature_focus", "스펙 불릿과 제품 강점을 강조하는 상세페이지형 구성"],
           ["promo_banner", "가격·채널·CTA를 전면에 둔 배너형 구성"]],
          col_widths=[3.0 * cm, 12.6 * cm])
    bullets([
        "문구 줄바꿈은 픽셀 폭 추정 기반으로 처리해 한글/영문 폭 차이에 따른 넘침을 줄인다.",
        "긴 단어 또는 공백 없는 한글도 글자 단위 분할 폴백으로 SVG 영역 밖으로 나가지 않게 한다.",
        "PPTX 다운로드는 ui/ppt_export.py가 SVG와 선택 문구를 사용해 생성한다.",
        "이미지 job이 실패해도 사용자는 이미지 없는 포스터를 생성할 수 있어 작업 흐름이 막히지 않는다.",
    ])

    # ── 11. 데이터와 스키마 ──
    story.append(PageBreak())
    para("11. 데이터 모델과 스키마", h1)
    para("요청 경계는 backend/schemas.py의 Pydantic 모델이 담당한다. 길이, 패턴, 범위를 API 입구에서 제한하고, AI 프롬프트 경로에서는 sanitize_user_text로 "
         "제어문자 제거와 공백 정규화를 다시 수행한다.")
    table([["스키마", "주요 필드와 제약"],
           ["KeyboardRenderRequest", "product_name 80자, layout, case/keycap/accent 색, case_finish, plate_material, switch_stem, keycap_profile, mount_type"],
           ["DeskSetupRenderRequest", "assets, desk_width 100~200cm, desk_depth 50~90cm, monitor_size 24/27/32, monitor_arm_style"],
           ["AdContentRequest", "상품/타깃/상세 설명, shot_type, image_ratio, reference image b64 최대 12MB, engine, selected_copy"],
           ["UploadedModelRequest", "filename path separator 금지, content_base64 최대 120MB"],
           ["SignupRequest", "username 영문/숫자 3~32자, password 8~128자, signup_code 1~64자"]],
          col_widths=[4.0 * cm, 11.6 * cm])
    para("데이터 파일은 data/layouts의 배열 JSON, data/desk_assets.json의 소품 카탈로그, data/drawings의 플레이트 도면 JSON으로 나뉜다. 런타임 산출물은 "
         "static/models, static/posters, data/runtime 아래에 생성된다. job store는 append-only jsonl을 replay해 최신 snapshot을 만들고, cache는 "
         "text/image 디렉터리에 SHA256 json으로 저장한다.")

    para("캐시와 영속성", h2)
    bullets([
        "텍스트 캐시 키는 상품명, 타깃, 소구점, 배열, 색상, 재질, provider, model, policy version을 포함한다.",
        "이미지 캐시 키는 image_prompt, workflow content hash, 해상도, ControlNet 설정, reference hash, shot_type, denoise/steps를 포함한다.",
        "이미지 캐시는 image_b64/image_b64s 바이너리를 저장하지 않고 job metadata만 저장해 디스크 증가를 줄인다.",
        "TTL과 최대 entry 수로 만료/축출하며, 읽기 시 mtime을 갱신해 LRU 기준으로 사용한다.",
    ])

    # ── 12. 보안 ──
    story.append(PageBreak())
    para("12. 인증과 보안", h1)
    para("보안 설계는 docs/security.md의 위협 모델을 따른다. 가장 중요한 위험은 secret 노출, 외부 무인증 접근, prompt injection, 업로드/path traversal, "
         "모델 워커 포트 노출이다.")
    table([["영역", "구현"],
           ["로그인/회원가입", "가입 코드 기반 회원가입, PBKDF2 200,000 iterations, salt 저장, 로그인 5회 실패 시 60초 잠금"],
           ["세션", "서버 메모리 세션 토큰, 12시간 TTL, /auth/session으로 새로고침 복원, logout 토큰 무효화"],
           ["사용자 저장소", "data/runtime/users.json에 평문 비밀번호 없이 salt+PBKDF2 해시 저장, 원자적 쓰기와 0600 권한"],
           ["시크릿 마스킹", "SENSITIVE_ENV_KEYS와 token-shaped regex 기반 로그 필터, /health와 /security/config는 set/missing만 노출"],
           ["프롬프트 인젝션", "한/영 우회 패턴을 flag로 기록하고 system prompt에서 내부 경로/키/시스템 프롬프트 노출 금지"],
           ["파일/경로", "업로드 파일명 패턴 제한, model_url은 basename만 신뢰, workflow 이름은 안전한 stem만 허용"],
           ["외부 노출", "Streamlit/FastAPI/ComfyUI/Ollama/HyperCLOVA 워커는 loopback, 외부는 nginx 8443 + Basic Auth 경유"]],
          col_widths=[3.4 * cm, 12.2 * cm])
    para("최종 로그인 업데이트에서는 Streamlit custom component iframe에서 top-level 자동 이동이 브라우저 정책에 막히는 문제를 확인하고, 로그인 UI가 더 이상 "
         "/auth/cookie 리다이렉트 URL을 표시하거나 자동 이동하지 않도록 정리했다. 결과적으로 로그인 성공 후 사용자는 같은 URL에서 메인 화면으로 전환된다.")

    # ── 13. 배포 ──
    story.append(PageBreak())
    para("13. 배포와 운영", h1)
    para("앱 티어는 FastAPI backend와 Streamlit frontend를 CPU-only Docker 컨테이너로 띄울 수 있다. GPU 작업은 컨테이너에 넣지 않고 host에서 관리되는 "
         "ComfyUI/Ollama/HyperCLOVA 워커를 HTTP로 호출한다. 단일 VM 구조에서는 docker-compose가 host network를 사용해 127.0.0.1 포트로 그대로 연결한다.")
    table([["서비스", "포트/위치", "운영 정책"],
           ["nginx", "외부 8443", "TLS + Basic Auth, Streamlit 8501로 WebSocket proxy"],
           ["Streamlit", "127.0.0.1:8501", "프론트 UI. 직접 외부 노출하지 않음"],
           ["FastAPI", "127.0.0.1:8010", "API backend. Streamlit이 내부 REST 호출"],
           ["ComfyUI", "127.0.0.1:8188", "FLUX/ControlNet 이미지 워커. systemd 또는 host 관리"],
           ["Ollama/local LLM", "127.0.0.1:11434", "선택적 로컬 텍스트 provider"],
           ["HyperCLOVA SEED/Omni", "11501/11601/11602", "텍스트, vision, image 워커. 단일 GPU와 VRAM 경합"]],
          col_widths=[3.2 * cm, 4.0 * cm, 8.4 * cm])
    para("GPU_WORKER_MODE는 always_on, on_demand, exclusive 세 모드를 제공한다. Docker 앱 티어에서는 워커를 컨테이너가 직접 기동하지 않으므로 always_on을 "
         "유지하고, host에서 워커를 관리하는 구성이 권장된다. 단일 실행 환경에서 exclusive를 쓰면 텍스트/이미지 워커 전환 시 경쟁 워커를 내려 VRAM을 확보한다.")
    para("운영 체크포인트", h2)
    bullets([
        ".env와 data/runtime 하위 파일 권한이 0600인지 확인한다.",
        "8010, 8501, 8188, 11434, 11501, 11601, 11602 포트가 loopback에만 바인딩되어 있는지 확인한다.",
        "코드 변경 후 라이브 검증 전에는 start.sh --restart 또는 docker compose up -d --build로 앱을 재기동한다.",
        "ComfyUI workflow, ControlNet 모델, LoRA, denoise/steps 변경은 이미지 cache key에 반영되므로 force_regen과 함께 검증한다.",
        "단일 L4에서는 FLUX+ControlNet과 Omni image가 동시에 상주하지 않도록 워커 모드를 관리한다.",
    ])

    # ── 14. 품질 보증 ──
    story.append(PageBreak())
    para("14. 품질 보증과 테스트", h1)
    para("테스트는 API/스키마/인증/렌더링/ComfyUI workflow/포스터 SVG/품질 게이트/캐시/QA 회귀를 폭넓게 포함한다. 시각 품질 자체는 단위 테스트로 완전히 "
         "판정하기 어려우므로, 라이브 ComfyUI 생성 결과를 이미지로 확인하는 절차를 병행한다.")
    table([["테스트 영역", "예시"],
           ["인증", "test_auth_login.py, test_auth_signup.py — 세션, 회원가입 코드, 잠금, 사용자 저장"],
           ["렌더링", "test_renderer_keyboard_silhouette.py, test_renderer_monitor_arm.py, test_setup_composition_raster.py, test_controlnet_depth.py"],
           ["AI/워크플로", "test_comfyui_workflow.py, test_hyperclova_job_queue.py, test_engine_selection_qa.py, test_llm_retry.py"],
           ["포스터/UI", "test_poster_svg_cta.py, test_ppt_export.py, test_ui_api.py, test_copy_prompt_qa.py"],
           ["품질/캐시", "test_quality_gate.py, test_result_cache.py, test_image_jobs_qa.py"],
           ["업로드/데이터", "test_cad_glb.py, test_drawing_converter.py, test_filenames.py"]],
          col_widths=[3.3 * cm, 12.3 * cm])
    para("quality_gate는 GPU/ML stack 없이도 동작하는 가벼운 신호를 제공한다. PNG/JPEG header에서 해상도와 ratio를 확인하고, PIL/numpy가 있으면 edge density와 "
         "frame fill로 sparse/front-elevation 구도를 잡는다. best-of-N에서는 중앙 band에서 액센트 색에 가까운 픽셀 비중을 비교해 색 충실도가 높은 컷을 "
         "대표 image_b64로 승격한다.")
    bullets([
        "배열 의미 정확도, 넘패드 오검출, 키캡 melt 같은 고차 의미 실패는 현재 lightweight gate만으로 완전히 잡지 못한다.",
        "따라서 최종 품질 검증은 단위 테스트 + 실제 ComfyUI 생성 결과 눈 검증 + 향후 VLM judge 도입의 3단계로 보는 것이 맞다.",
        "회귀 테스트 실행 기준 명령은 conda run -n sprint_high pytest이며, 보고서 생성 자체는 tools/build_presentation.py로 검증한다.",
    ])

    # ── 15. 성과와 한계 ──
    story.append(PageBreak())
    para("15. 성과, 한계, 향후 과제", h1)
    para("DeskAd의 현재 성과는 ‘광고 이미지를 빠르게 만든다’보다 넓다. 판매자가 입력한 상품 스펙을 렌더링, 프롬프트, ControlNet, 포스터까지 하나의 "
         "데이터 흐름으로 유지하는 구조를 만들었다. 이 구조 덕분에 생성 이미지가 실제 제품과 달라지는 문제를 시스템적으로 줄일 수 있다.")
    table([["성과", "설명"],
           ["정확한 기준 데이터", "레이아웃 JSON과 cm 단위 GLB가 광고 생성의 기준이 되어 프롬프트만 쓰는 방식보다 통제가 강함"],
           ["작업 완결성", "문구 후보, 이미지 job, 포스터 SVG/PPTX까지 이어져 실제 제출/판매 소재에 가까운 산출물을 제공"],
           ["실패 복구", "not_configured/failed 상태, 이미지 없는 포스터, fallback copy/SVG, force_regen으로 데모가 중단되지 않음"],
           ["운영성", "Docker 앱 티어, host GPU 워커, 캐시, job store, 워커 idle/release로 반복 운영 기반 확보"],
           ["보안성", "로그/응답 시크릿 마스킹, 회원가입 코드, 세션, 파일 권한, pre-commit scan, loopback 바인딩"]],
          col_widths=[3.2 * cm, 12.4 * cm])
    para("남은 한계", h2)
    bullets([
        "ControlNet depth는 구조를 고정하지만 색은 고정하지 못하므로 프롬프트와 best-of-N에 의존한다.",
        "lightweight quality gate는 구도와 색의 일부 신호만 보며, 키 배열의 의미적 정확도는 VLM 기반 판정이 필요하다.",
        "ComfyUI/HyperCLOVA 워커는 단일 L4 VRAM을 공유하므로, 고해상도·다중 batch·여러 동시 사용자는 별도 큐/자원 스케줄러가 필요하다.",
        "현재 3D 렌더러는 절차적 primitive 기반이라 실제 제품의 곡면/나사/키캡 legends를 완전한 CAD 수준으로 재현하지는 않는다.",
        "PPTX/포스터 템플릿은 실사용을 위해 더 많은 채널별 변형과 긴 텍스트 검수 케이스가 필요하다.",
    ])
    para("향후 과제", h2)
    bullets([
        "VLM judge를 도입해 배열, 넘패드, 키 열 직선성, 키캡 melt, 소품 중복을 의미적으로 평가한다.",
        "ControlNet strength/end_percent와 best-of-N을 사용자 프리셋으로 묶어 ‘정확도 우선/사진감 우선’ 모드를 제공한다.",
        "이미지 생성 job queue를 사용자/세션 단위로 분리하고 동시 사용자 제한, 우선순위, 재시도 정책을 명시한다.",
        "실제 판매자의 CAD/사진 레퍼런스가 있을 때 GLB proxy가 아니라 정확한 모델 변환 파이프라인을 강화한다.",
        "포스터 템플릿을 채널별 규격으로 확장하고, 상세페이지용 세로 긴 이미지/카피 블록 산출물을 추가한다.",
    ])

    para("부록: 주요 코드 경로", h1)
    table([["파일", "보고서에서 사용한 근거"],
           ["docs/architecture.md", "시스템 구조, depth-ControlNet, GPU 워커, 품질 검증 방침"],
           ["docs/schema.md", "API 엔드포인트와 Pydantic 요청/응답 모델"],
           ["docs/security.md", "위협 모델, 시크릿 위생, 입력 검증, nginx/basic auth, 파일 권한"],
           ["docs/deploy.md", "Docker 앱 티어, host GPU 워커 연결, 컷오버 절차"],
           ["backend/renderer.py", "GLB 생성, composition raster, OSMesa depth PNG"],
           ["backend/ai.py", "카피/이미지 프롬프트, workflow 선택, ComfyUI job, grid_three, SVG 포스터"],
           ["backend/runtime_workers.py", "GPU_WORKER_MODE, 워커 예열/해제/idle reap"],
           ["backend/auth.py, backend/user_store.py", "세션 인증, 회원가입, PBKDF2 사용자 저장"],
           ["ui/steps.py, ui/ad_content.py", "4단계 UI, 작업 상태 카드, 자동 폴링, 문구 후보 선택"],
           ["worklogs/성현/2026-06-10~15report.md", "최종 업데이트 UX/인증/로딩/오류 처리 작업 내역"]],
          col_widths=[5.0 * cm, 10.6 * cm])

    doc = SimpleDocTemplate(str(PDF_OUT), pagesize=A4,
                            leftMargin=2.2 * cm, rightMargin=2.2 * cm,
                            topMargin=2.0 * cm, bottomMargin=1.8 * cm,
                            title="DeskAd AI Studio 프로젝트 보고서", author="DeskAd")

    def footer(canvas, d):
        canvas.saveState()
        canvas.setFont("Nanum", 8)
        canvas.setFillColor(colors.HexColor(MUTED))
        canvas.drawString(2.2 * cm, 1.1 * cm, "DeskAd AI Studio · 프로젝트 보고서")
        canvas.drawRightString(A4[0] - 2.2 * cm, 1.1 * cm, "%d" % d.page)
        canvas.restoreState()

    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    print("PDF ->", PDF_OUT)


# ============================================================================
# 2) PPTX 슬라이드 (python-pptx)
# ============================================================================
def build_pptx() -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    def rgb(h):
        return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))

    KOR = "맑은 고딕"  # PowerPoint 한글 렌더용(없으면 시스템 한글 폰트로 대체)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def slide():
        return prs.slides.add_slide(blank)

    def rect(s, x, y, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(color)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def textbox(s, x, y, w, h, lines, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font=KOR):
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        if isinstance(lines, str):
            lines = [lines]
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if isinstance(ln, tuple):
                text, lvl, sz, bd = ln
            else:
                text, lvl, sz, bd = ln, 0, size, bold
            p.level = lvl
            run = p.add_run()
            run.text = text
            run.font.size = Pt(sz)
            run.font.bold = bd
            run.font.name = font
            run.font.color.rgb = rgb(color if not isinstance(ln, tuple) else INK)
        return tb

    def bullets_box(s, x, y, w, h, items, size=16):
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            lvl = 0
            txt = it
            if isinstance(it, tuple):
                txt, lvl = it
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = lvl
            run = p.add_run()
            run.text = ("• " if lvl == 0 else "– ") + txt
            run.font.size = Pt(size - lvl * 2)
            run.font.name = KOR
            run.font.color.rgb = rgb(INK if lvl == 0 else MUTED)
            p.space_after = Pt(6)
        return tb

    def pic_fit(s, name, x, y, max_w, max_h):
        sz = img_size(name)
        if not sz:
            return
        w, h = sz
        scale = min(max_w / w, max_h / h)
        dw, dh = int(w * scale), int(h * scale)
        s.shapes.add_picture(str(IMG / name), x + (max_w - dw) // 2, y + (max_h - dh) // 2, dw, dh)

    def header(s, title, idx=None):
        rect(s, 0, 0, SW, Inches(1.15), INK)
        rect(s, 0, Inches(1.15), SW, Pt(4), ACCENT)
        textbox(s, Inches(0.55), Inches(0.18), SW - Inches(1.1), Inches(0.85),
                title, size=26, color="#ffffff", bold=True, anchor=MSO_ANCHOR.MIDDLE)
        if idx:
            textbox(s, SW - Inches(1.4), Inches(0.18), Inches(0.9), Inches(0.85),
                    idx, size=13, color="#c8c1b2", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # ── 1. 표지 ──
    s = slide()
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, Inches(4.05), SW, Pt(5), ACCENT)
    textbox(s, Inches(0.8), Inches(2.2), SW - Inches(1.6), Inches(1.4),
            "DeskAd AI Studio", size=52, color="#ffffff", bold=True)
    textbox(s, Inches(0.85), Inches(4.25), SW - Inches(1.6), Inches(1.6),
            ["3D 데스크 셋업 미리보기 + 배열 충실도 기반 광고 콘텐츠 자동 생성",
             "소상공인 커스텀 키보드 · 데스크테리어 판매자용 · 2026-06-17"],
            size=18, color="#c8c1b2")

    # ── 2. 목차 ──
    s = slide()
    header(s, "목차", "01")
    bullets_box(s, Inches(0.7), Inches(1.5), SW - Inches(1.4), Inches(5.6), [
        "배경과 문제 — ‘AI 광고가 실제 제품과 다르다’",
        "시스템 개요 — UI · 백엔드 · 2트랙 엔진",
        "사용 흐름 — 4단계 위저드",
        "핵심 기술 — 배열 충실도(depth-ControlNet)",
        "그리드 광고 — 컷별 시점 분리 + 순차 생성",
        "보안 · 품질 · 성과",
    ], size=20)

    # ── 3. 문제 ──
    s = slide()
    header(s, "배경과 문제 정의", "02")
    bullets_box(s, Inches(0.7), Inches(1.5), Inches(6.4), Inches(5.4), [
        "소상공인은 광고 촬영·디자인에 비용·시간을 들이기 어렵다.",
        "범용 생성 모델은 빠르지만 정밀 제품에서 배열이 틀린다:",
        ("65% 키보드를 풀사이즈로 생성", 1),
        ("행 물결침 · 키캡 뭉개짐(melt)", 1),
        "→ ‘그럴듯하지만 실제 상품과 다른’ 이미지 = 광고로 못 씀.",
        "DeskAd: 실제 스펙의 3D를 구조로 강제 주입해 해결.",
    ], size=17)
    pic_fit(s, "ad_hero.png", Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.0))

    # ── 4. 시스템 개요 ──
    s = slide()
    header(s, "시스템 개요", "03")
    bullets_box(s, Inches(0.7), Inches(1.45), Inches(6.5), Inches(5.6), [
        "브라우저 → nginx(Basic Auth) → Streamlit UI → FastAPI 백엔드",
        "백엔드 = 단일 오케스트레이터(검증·프롬프트·잡·캐시·워커)",
        "이미지/텍스트 2트랙:",
        ("local: 로컬 LLM + ComfyUI/FLUX (정확도 우선, 키 불필요)", 1),
        ("openai: OpenAI 텍스트 + 이미지 (키 필요)", 1),
        "GPU 워커는 단일 L4(24GB) 공유 — GPU_WORKER_MODE로 제어",
        "엔진/키 없으면 템플릿·SVG로 안전 폴백",
    ], size=16)
    pic_fit(s, "ui_3d_preview.png", Inches(7.45), Inches(1.7), Inches(5.3), Inches(4.8))

    # ── 5. 4단계 ──
    s = slide()
    header(s, "사용 흐름 — 4단계 위저드", "04")
    bullets_box(s, Inches(0.7), Inches(1.45), Inches(6.4), Inches(5.6), [
        "① 상품 정보 — 모델명·레이아웃, STEP/GLB 업로드",
        "② 도면 미리보기 — 레이아웃 → SVG 탑뷰",
        "③ 3D 셋업 — 색·재질·책상/모니터/소품 → GLB(1 unit=1cm)",
        "④ 광고 생성 — 문구 + 이미지 + SVG/PPTX 포스터",
        ("MX 표준 간격 19.05mm·실측 footprint → 충실도의 토대", 1),
    ], size=17)
    pic_fit(s, "setup_3d_render.png", Inches(7.4), Inches(1.6), Inches(5.4), Inches(5.2))

    # ── 6. 핵심 기술 ──
    s = slide()
    header(s, "핵심 기술 — 배열 충실도 (depth-ControlNet)", "05")
    textbox(s, Inches(0.7), Inches(1.3), SW - Inches(1.4), Inches(0.9),
            "셋업 GLB →(CPU 헤드리스 렌더) depth →  FLUX depth-ControlNet → 사진 광고  ·  구조와 외관을 분리",
            size=16, color=ACCENT, bold=True)
    y = Inches(2.3)
    for i, (name, capt) in enumerate([("composition_perspective.png", "① 셋업(구도 맵)"),
                                      ("depth_hero.png", "② depth 입력"),
                                      ("ad_hero.png", "③ 최종 사진 광고")]):
        x = Inches(0.7 + i * 4.25)
        pic_fit(s, name, x, y, Inches(3.9), Inches(3.6))
        textbox(s, x, Inches(5.95), Inches(3.9), Inches(0.5), capt, size=14,
                color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    textbox(s, Inches(0.7), Inches(6.55), SW - Inches(1.4), Inches(0.7),
            "배열=depth · 주색=프롬프트 · 액센트=best-of-N (3단 직교) · strength 0.5 = 사진+정확 배열 스위트스팟",
            size=14, color=INK)

    # ── 7. 컷별 시점 ──
    s = slide()
    header(s, "그리드 광고 — 컷별 시점 분리", "06")
    for i, (name, capt) in enumerate([("depth_hero.png", "depth: hero(높은 3/4)"),
                                      ("depth_eye_level.png", "depth: eye_level(낮은 수평)"),
                                      ("depth_legacy.png", "(이전) 고정 각도 — 겹침")]):
        x = Inches(0.7 + i * 4.25)
        pic_fit(s, name, x, Inches(1.4), Inches(3.9), Inches(2.5))
        textbox(s, x, Inches(3.85), Inches(3.9), Inches(0.4), capt, size=12,
                color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    for i, (name, capt) in enumerate([("ad_hero.png", "최종 hero"),
                                      ("ad_eye_level.png", "최종 eye_level"),
                                      ("ad_detail_macro.png", "최종 detail_macro")]):
        x = Inches(0.7 + i * 4.25)
        pic_fit(s, name, x, Inches(4.35), Inches(3.9), Inches(2.4))
        textbox(s, x, Inches(6.75), Inches(3.9), Inches(0.4), capt, size=12,
                color=MUTED, align=PP_ALIGN.CENTER, bold=True)

    # ── 8. 순차 생성 ──
    s = slide()
    header(s, "그리드 — 순차 생성(안정성)", "07")
    bullets_box(s, Inches(0.7), Inches(1.5), SW - Inches(1.4), Inches(5.5), [
        "문제: 3컷을 한꺼번에 큐잉 → 단일 L4 VRAM 피크 겹침 → 서버 다운 위험",
        "해법: 첫 컷만 제출 → 폴링이 완료 확인 후에야 다음 컷 제출",
        ("ComfyUI 큐에 우리 컷이 항상 1개만 존재", 1),
        ("컷별 depth는 유니크 파일명으로 업로드(overwrite 클로버 방지)", 1),
        "라이브 검증: in-flight=1 유지, 크래시 없이 3컷 완주",
        "회귀 테스트 263개 통과 + 라이브 눈 검증",
    ], size=18)

    # ── 9. 보안·품질 ──
    s = slide()
    header(s, "보안 · 품질 · 성과", "08")
    bullets_box(s, Inches(0.7), Inches(1.45), Inches(6.4), Inches(5.6), [
        "보안(docs/security.md):",
        ("시크릿 마스킹 + pre-commit 스캔", 1),
        ("nginx Basic Auth + 세션 인증, 워커 포트 비공개", 1),
        ("프롬프트 인젝션 탐지 + Pydantic 입력 검증", 1),
        "품질:",
        ("회귀 263개 통과 · 라이브 시각 검증", 1),
        ("strength 0.5 = 사진+정확 배열 동시 달성", 1),
    ], size=16)
    pic_fit(s, "ad_eye_level.png", Inches(7.5), Inches(1.6), Inches(5.2), Inches(5.0))

    # ── 10. 마무리 ──
    s = slide()
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, Inches(3.7), SW, Pt(5), ACCENT)
    textbox(s, Inches(0.8), Inches(2.6), SW - Inches(1.6), Inches(1.2),
            "정확한 3D → 정확한 광고", size=40, color="#ffffff", bold=True)
    textbox(s, Inches(0.85), Inches(3.95), SW - Inches(1.6), Inches(1.2),
            ["DeskAd AI Studio — 배열 충실도로 ‘쓸 수 있는’ 생성형 광고",
             "감사합니다."], size=20, color="#c8c1b2")

    prs.save(str(PPTX_OUT))
    print("PPTX ->", PPTX_OUT, f"({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    PRES.mkdir(parents=True, exist_ok=True)
    build_pdf()
    build_pptx()
